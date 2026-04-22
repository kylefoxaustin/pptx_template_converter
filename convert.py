#!/usr/bin/env python3
"""Apply a corporate PowerPoint template's theme/master/layouts to an existing deck.

Strategy A — "theme swap": use the template as the base, port source slides onto
the template's Blank layout. Source's absolutely-positioned shapes are preserved
exactly; the template's theme (colors, fonts) and layout branding take over.

Usage:
    python convert.py --input INPUT.pptx --template TEMPLATE.pptx --output OUTPUT.pptx
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.parts.image import Image, ImagePart

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
SRGB_RE = re.compile(r'srgbClr\s+val="([0-9A-Fa-f]{6})"')

# spTree children that must stay put when we wipe a slide's default shapes.
GROUP_FRAME_TAGS = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}

THEME_COLOR_NAMES = {
    "dk1", "lt1", "dk2", "lt2", "tx1", "tx2", "bg1", "bg2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
}

SCHEME_TO_MSO = {
    "dk1": MSO_THEME_COLOR.DARK_1,
    "lt1": MSO_THEME_COLOR.LIGHT_1,
    "dk2": MSO_THEME_COLOR.DARK_2,
    "lt2": MSO_THEME_COLOR.LIGHT_2,
    "tx1": MSO_THEME_COLOR.TEXT_1,
    "tx2": MSO_THEME_COLOR.TEXT_2,
    "bg1": MSO_THEME_COLOR.BACKGROUND_1,
    "bg2": MSO_THEME_COLOR.BACKGROUND_2,
    "accent1": MSO_THEME_COLOR.ACCENT_1,
    "accent2": MSO_THEME_COLOR.ACCENT_2,
    "accent3": MSO_THEME_COLOR.ACCENT_3,
    "accent4": MSO_THEME_COLOR.ACCENT_4,
    "accent5": MSO_THEME_COLOR.ACCENT_5,
    "accent6": MSO_THEME_COLOR.ACCENT_6,
    "hlink": MSO_THEME_COLOR.HYPERLINK,
    "folHlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
}
HEX_RE = re.compile(r"^#?[0-9A-Fa-f]{6}$")

# Title detection: a shape has to span at least this fraction of the slide width
# to be considered a title candidate. Titles typically run nearly the full width;
# decorative stat boxes / column cards span ~1/3.
TITLE_MIN_WIDTH_FRACTION = 0.50

# Fallback (no shape meets the width bar): require the biggest-font shape to
# exceed the next-largest by at least this factor before we treat it as a title.
TITLE_SIZE_RATIO = 1.15


def convert(src_path: Path, tpl_path: Path, out_path: Path,
            color_map_path: Path | None = None,
            title_color: str = "auto") -> None:
    src = Presentation(src_path)
    dst = Presentation(tpl_path)

    if (src.slide_width, src.slide_height) != (dst.slide_width, dst.slide_height):
        # Ignore sub-inch rounding differences that show up between decks authored in
        # different tools — they don't warrant a warning.
        width_delta = abs(src.slide_width - dst.slide_width)
        height_delta = abs(src.slide_height - dst.slide_height)
        if width_delta > 9144 or height_delta > 9144:  # > 0.01"
            print(
                f"[warn] canvas size differs: source {src.slide_width}x{src.slide_height}, "
                f"template {dst.slide_width}x{dst.slide_height}. Keeping source dimensions; "
                f"template branding designed for the other ratio may look stretched.",
                file=sys.stderr,
            )
        dst.slide_width = src.slide_width
        dst.slide_height = src.slide_height

    blank_layout = _find_blank_layout(dst)
    print(f"[info] using template layout: {blank_layout.name!r}")

    theme_colors = _theme_colors(dst)
    color_map = _load_color_map(color_map_path, theme_colors) if color_map_path else {}
    if color_map_path:
        print(f"[info] loaded {len(color_map)} color mappings from {color_map_path}")

    title_color_spec = _resolve_title_color(title_color, dst)

    _delete_all_slides(dst)

    per_slide_report: list[tuple[int, dict[str, set[str]]]] = []
    for idx, src_slide in enumerate(src.slides):
        new_slide = dst.slides.add_slide(blank_layout)
        _strip_default_shapes(new_slide)
        rid_map = _clone_slide_rels(src_slide, new_slide, idx)
        _copy_shapes(src_slide, new_slide, rid_map, color_map)
        per_slide_report.append((idx, _categorize_colors(src_slide, color_map)))

    title_log = _recolor_titles(dst, title_color_spec) if title_color_spec else []

    out_path.parent.mkdir(parents=True, exist_ok=True)
    dst.save(out_path)
    print(f"[done] wrote {out_path}")
    _print_color_report(per_slide_report, theme_colors, bool(color_map))
    _print_title_report(title_log, title_color_spec)


def _find_blank_layout(pres):
    for layout in pres.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return pres.slide_layouts[0]


def _delete_all_slides(pres):
    """Remove every slide from `pres` — drop the sldId entries AND their relationships."""
    sld_id_lst = pres.slides._sldIdLst
    r_id_attr = f"{{{R_NS}}}id"
    for sld_id in list(sld_id_lst):
        rid = sld_id.get(r_id_attr)
        pres.part.drop_rel(rid)
        sld_id_lst.remove(sld_id)


def _strip_default_shapes(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        if child.tag not in GROUP_FRAME_TAGS:
            sp_tree.remove(child)


def _clone_slide_rels(src_slide, dst_slide, slide_idx):
    """Port images and hyperlinks from src_slide onto dst_slide's part.

    Returns a map of src-side rId -> dst-side rId so we can remap r:embed / r:link / r:id
    references on the copied shape XML.
    """
    rid_map: dict[str, str] = {}
    for rel in src_slide.part.rels.values():
        # dst_slide already has the correct layout relationship; don't overwrite.
        if rel.reltype.endswith("/slideLayout"):
            continue

        if rel.is_external:
            new_rid = dst_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
            rid_map[rel.rId] = new_rid
            continue

        src_part = rel.target_part
        if isinstance(src_part, ImagePart):
            img = Image.from_blob(src_part.blob, src_part.desc)
            new_part = ImagePart.new(dst_slide.part.package, img)
            new_rid = dst_slide.part.relate_to(new_part, rel.reltype)
            rid_map[rel.rId] = new_rid
        else:
            print(
                f"[warn] slide {idx_for_log(slide_idx)}: skipping related part "
                f"{src_part.partname} (reltype={rel.reltype}). Visuals tied to it may be broken.",
                file=sys.stderr,
            )
    return rid_map


def idx_for_log(i: int) -> str:
    return f"{i + 1}"


def _copy_shapes(src_slide, dst_slide, rid_map, color_map):
    dst_tree = dst_slide.shapes._spTree
    r_attrs = (f"{{{R_NS}}}embed", f"{{{R_NS}}}link", f"{{{R_NS}}}id")
    for src_el in list(src_slide.shapes._spTree):
        if src_el.tag in GROUP_FRAME_TAGS:
            continue
        new_el = deepcopy(src_el)
        for el in new_el.iter():
            for attr in r_attrs:
                old = el.attrib.get(attr)
                if old is not None and old in rid_map:
                    el.attrib[attr] = rid_map[old]
        if color_map:
            _remap_colors(new_el, color_map)
        dst_tree.append(new_el)


def _remap_colors(root, color_map):
    """Rewrite matching <a:srgbClr val="X"/> elements per color_map.

    A mapping target can be either a template theme-slot name (rendered as
    <a:schemeClr val="accent1"/>) or a hex literal (stays as <a:srgbClr>).
    Preserves any shade/tint/alpha child modifiers on the original color element.
    """
    srgb_tag = f"{{{A_NS}}}srgbClr"
    scheme_tag = f"{{{A_NS}}}schemeClr"
    for old_el in list(root.iter(srgb_tag)):
        val = (old_el.get("val") or "").upper()
        target = color_map.get(val)
        if target is None:
            continue
        is_theme = not target.startswith("#")
        new_el = etree.Element(scheme_tag if is_theme else srgb_tag)
        new_el.set("val", target.lstrip("#") if not is_theme else target)
        for child in list(old_el):
            new_el.append(child)
        old_el.addnext(new_el)
        old_el.getparent().remove(old_el)


def _load_color_map(path: Path, theme_colors: dict[str, str]) -> dict[str, str]:
    """Load a color-map JSON file and validate entries.

    The file is { "description": "...", "map": { "<srcHex>": "<target>", ... } }.
    Targets are either a theme-slot name (dk1/lt1/dk2/lt2/accent1..accent6/hlink/folHlink)
    or a '#RRGGBB' literal. Unknown theme names cause a fatal error so typos don't
    silently no-op.
    """
    data = json.loads(path.read_text())
    raw = data.get("map", {}) if isinstance(data, dict) else {}
    normalized: dict[str, str] = {}
    for src, target in raw.items():
        if not isinstance(src, str) or not HEX_RE.match(src):
            raise SystemExit(f"color-map key not a 6-hex string: {src!r}")
        src_key = src.lstrip("#").upper()
        if not isinstance(target, str):
            raise SystemExit(f"color-map value for {src!r} is not a string: {target!r}")
        if target.startswith("#"):
            if not HEX_RE.match(target):
                raise SystemExit(f"color-map hex target for {src!r} invalid: {target!r}")
            normalized[src_key] = "#" + target.lstrip("#").upper()
        elif target in THEME_COLOR_NAMES:
            if target in ("accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                          "dk1", "lt1", "dk2", "lt2") and target not in theme_colors:
                print(f"[warn] color-map: template theme has no {target!r}; mapping of "
                      f"{src!r} will still emit a schemeClr reference (PowerPoint will resolve).",
                      file=sys.stderr)
            normalized[src_key] = target
        else:
            raise SystemExit(
                f"color-map target for {src!r} not recognized: {target!r}. "
                f"Expected a theme name ({', '.join(sorted(THEME_COLOR_NAMES))}) or '#RRGGBB'."
            )
    return normalized


def _categorize_colors(slide, color_map) -> dict[str, set[str]]:
    """Return per-slide color buckets: 'remapped' (source hexes that were hit by color_map)
    and 'unmapped' (source hexes with no mapping — still hardcoded in output)."""
    xml = etree.tostring(slide._element, encoding="unicode")
    all_colors = {m.group(1).upper() for m in SRGB_RE.finditer(xml)}
    return {
        "remapped": {c for c in all_colors if c in color_map},
        "unmapped": {c for c in all_colors if c not in color_map},
    }


def _theme_colors(pres) -> dict[str, str]:
    master = pres.slide_masters[0]
    theme_part = next(
        (r.target_part for r in master.part.rels.values() if r.reltype.endswith("/theme")),
        None,
    )
    if theme_part is None:
        return {}
    root = etree.fromstring(theme_part.blob)
    out: dict[str, str] = {}
    for name in ("dk1", "lt1", "dk2", "lt2",
                 "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"):
        el = root.find(f".//{{{A_NS}}}{name}")
        if el is None:
            continue
        srgb = el.find(f"{{{A_NS}}}srgbClr")
        if srgb is not None:
            out[name] = srgb.get("val").upper()
    return out


def _print_color_report(per_slide_report, theme_colors, used_color_map: bool):
    print()
    print("=== Color report ===")
    if theme_colors:
        print("Template theme palette:")
        for name, hx in theme_colors.items():
            print(f"  {name:8s} #{hx}")
    print()
    # Aggregate across slides so the user sees the full set of unmapped colors to address.
    all_remapped: set[str] = set()
    all_unmapped: set[str] = set()
    for _, buckets in per_slide_report:
        all_remapped |= buckets["remapped"]
        all_unmapped |= buckets["unmapped"]

    if used_color_map:
        print(f"Source colors remapped ({len(all_remapped)} unique):")
        if all_remapped:
            print("  " + ", ".join("#" + c for c in sorted(all_remapped)))
        else:
            print("  (none — the mapping had no hits on this source)")
        print()
        print(f"Source colors still hardcoded — add these to your color map to re-theme them "
              f"({len(all_unmapped)} unique):")
        if all_unmapped:
            print("  " + ", ".join("#" + c for c in sorted(all_unmapped)))
        else:
            print("  (none — mapping covered everything)")
    else:
        print("Source colors hardcoded in output (no --color-map supplied):")
        print("  " + ", ".join("#" + c for c in sorted(all_unmapped)))


def _resolve_title_color(spec: str, pres):
    """Return a resolved title color or None (skip title recoloring).

    spec is one of:
      'auto'   — read from the template's master titleStyle; fallback accent1.
      'none'   — skip the title recoloring pass entirely.
      'dk1' .. — a theme color slot name.
      '#RRGGBB' — a hex literal.
    Return value: ('scheme', <slot>) or ('srgb', <HEX>) or None.
    """
    if spec == "none":
        return None
    if spec == "auto":
        resolved = _read_master_title_color(pres)
        if resolved is None:
            print("[info] could not read title color from template master; defaulting to accent1")
            return ("scheme", "accent1")
        print(f"[info] template title color detected: {resolved[0]}={resolved[1]}")
        return resolved
    if spec.startswith("#"):
        hex_ = spec.lstrip("#").upper()
        if not HEX_RE.match(hex_):
            raise SystemExit(f"Invalid --title-color hex: {spec}")
        return ("srgb", hex_)
    if spec in THEME_COLOR_NAMES:
        return ("scheme", spec)
    raise SystemExit(
        f"Invalid --title-color: {spec!r}. Use a theme slot name "
        f"({', '.join(sorted(THEME_COLOR_NAMES))}), '#RRGGBB', 'auto', or 'none'."
    )


def _read_master_title_color(pres):
    """Read the master's titleStyle level-1 solid fill, as a ('scheme'|'srgb', value) tuple."""
    master = pres.slide_masters[0]
    title_style = master.element.find(f".//{{{P_NS}}}titleStyle")
    if title_style is None:
        return None
    lvl1 = title_style.find(f"{{{A_NS}}}lvl1pPr")
    if lvl1 is None:
        return None
    def_rpr = lvl1.find(f"{{{A_NS}}}defRPr")
    if def_rpr is None:
        return None
    fill = def_rpr.find(f"{{{A_NS}}}solidFill")
    if fill is None:
        return None
    scheme = fill.find(f"{{{A_NS}}}schemeClr")
    if scheme is not None:
        val = scheme.get("val")
        if val in THEME_COLOR_NAMES:
            return ("scheme", val)
    srgb = fill.find(f"{{{A_NS}}}srgbClr")
    if srgb is not None:
        return ("srgb", srgb.get("val").upper())
    return None


def _detect_title_shape(slide, slide_width, slide_height):
    """Pick the shape that looks like a title.

    Primary heuristic: the *topmost shape that spans most of the slide width*
    is the title. Titles are authored full-width near the top; decorative stat
    cards / column boxes are narrower even when their font is larger.

    Fallback (no wide shape — unusual layouts): the overall largest-font shape,
    but only if it beats the next-largest by TITLE_SIZE_RATIO so we don't
    mislabel one of several evenly-sized body lines.
    """
    width_threshold = slide_width * TITLE_MIN_WIDTH_FRACTION
    wide_candidates: list[tuple[int, float, object]] = []  # (top, max_size, shape)
    all_candidates: list[tuple[float, int, object]] = []   # (max_size, top, shape)

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = (sh.text_frame.text or "").strip()
        if not text:
            continue  # empty textboxes would steal "topmost" from the real title
        max_size = 0.0
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                size = run.font.size
                if size is not None and size.pt > max_size:
                    max_size = size.pt
        top = sh.top if sh.top is not None else 0
        width = sh.width if sh.width is not None else 0
        all_candidates.append((max_size, top, sh))
        if width >= width_threshold:
            wide_candidates.append((top, max_size, sh))

    if wide_candidates:
        wide_candidates.sort(key=lambda c: (c[0], -c[1]))
        return wide_candidates[0][2]

    if not all_candidates:
        return None
    all_candidates.sort(key=lambda c: (-c[0], c[1]))
    best_size, _, best_shape = all_candidates[0]
    if best_size == 0:
        return None
    other_sizes = [c[0] for c in all_candidates[1:] if c[0] < best_size]
    if other_sizes and best_size < other_sizes[0] * TITLE_SIZE_RATIO:
        return None
    return best_shape


def _apply_title_color(shape, color_spec):
    kind, value = color_spec
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            color = run.font.color
            if kind == "scheme":
                color.theme_color = SCHEME_TO_MSO[value]
            else:
                color.rgb = RGBColor.from_string(value.lstrip("#"))


def _recolor_titles(pres, color_spec):
    """Apply color_spec to the detected title shape on every slide. Return log tuples."""
    log = []
    slide_width = pres.slide_width
    slide_height = pres.slide_height
    for idx, slide in enumerate(pres.slides):
        title_shape = _detect_title_shape(slide, slide_width, slide_height)
        if title_shape is None:
            continue
        _apply_title_color(title_shape, color_spec)
        snippet = (title_shape.text_frame.text.strip().replace("\n", " ") or "")[:60]
        log.append((idx + 1, snippet))
    return log


def _print_title_report(log, color_spec):
    print()
    print("=== Title recolor report ===")
    if color_spec is None:
        print("(skipped — title recoloring disabled)")
        return
    kind, value = color_spec
    target = f"{kind}={value}" if kind == "scheme" else f"#{value}"
    if not log:
        print(f"(no slides matched the title heuristic; font-size gaps too small)")
        print(f"Target color was: {target}")
        return
    print(f"Applied color {target} to {len(log)} slide title(s):")
    for slide_num, text in log:
        print(f"  slide {slide_num}: {text!r}")


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--input", required=True, type=Path, help="Source .pptx")
    ap.add_argument("--template", required=True, type=Path, help="Corporate template .pptx")
    ap.add_argument("--output", required=True, type=Path, help="Output merged .pptx")
    ap.add_argument("--color-map", type=Path, default=None,
                    help="Optional JSON file mapping source hex colors to template theme slots "
                         "or hex literals. See mappings/keyhole_to_corporate.json for an example.")
    ap.add_argument("--title-color", default="auto",
                    help="Color for detected slide titles: 'auto' (read from template master), "
                         "'none' (skip), a theme slot name (e.g. 'accent1'), or '#RRGGBB'.")
    args = ap.parse_args()
    for label, p in (("input", args.input), ("template", args.template)):
        if not p.exists():
            ap.error(f"{label} not found: {p}")
    if args.color_map is not None and not args.color_map.exists():
        ap.error(f"color-map not found: {args.color_map}")
    convert(args.input, args.template, args.output,
            color_map_path=args.color_map, title_color=args.title_color)


if __name__ == "__main__":
    main()
