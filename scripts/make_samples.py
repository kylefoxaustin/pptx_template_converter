#!/usr/bin/env python3
"""Generate the two synthetic sample fixtures.

    python scripts/make_samples.py

Writes:
    input/sample_input.pptx       — AI-generated-style deck
    template/sample_template.pptx — "Acme Corporation" branded template
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
INPUT_PATH = ROOT / "input" / "sample_input.pptx"
TEMPLATE_PATH = ROOT / "template" / "sample_template.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# --- sample_input: AI-generated-style deck ----------------------------------

NAVY = RGBColor(0x1F, 0x4E, 0x79)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)


def _set_run(tf, text, *, size, bold=False, color=DARK, align=None):
    tf.text = text
    para = tf.paragraphs[0]
    if align is not None:
        para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_input() -> None:
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    blank = p.slide_layouts[6]

    # Slide 1 — title slide
    s = p.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, p.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_GRAY
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), p.slide_width, Inches(2.5))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    title = s.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(12), Inches(1.2))
    _set_run(title.text_frame, "Quarterly Business Review", size=44, bold=True, color=WHITE)
    sub = s.shapes.add_textbox(Inches(0.75), Inches(4.0), Inches(12), Inches(0.6))
    _set_run(sub.text_frame, "Q4 2025", size=24, color=WHITE)

    # Slide 2 — agenda bullets
    s = p.slides.add_slide(blank)
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    _set_run(t.text_frame, "Agenda", size=32, bold=True, color=NAVY)
    body = s.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(12), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    items = [
        "Q4 revenue and growth highlights",
        "Customer acquisition and retention",
        "Product roadmap update",
        "Operational efficiency gains",
        "Outlook and 2026 priorities",
    ]
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = "•  " + item
        para.runs[0].font.size = Pt(22)
        para.runs[0].font.color.rgb = DARK
        para.space_after = Pt(12)

    # Slide 3 — table
    s = p.slides.add_slide(blank)
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    _set_run(t.text_frame, "Key Metrics", size=32, bold=True, color=NAVY)
    rows, cols = 4, 3
    tbl = s.shapes.add_table(rows, cols, Inches(1.0), Inches(1.6), Inches(11.3), Inches(3.5)).table
    headers = ["Metric", "Q3 2025", "Q4 2025"]
    data = [
        ["Revenue", "$12.4M", "$14.1M"],
        ["New customers", "82", "104"],
        ["Churn rate", "3.2%", "2.7%"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.size = Pt(16)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for ri, row in enumerate(data, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.color.rgb = DARK

    # Slide 4 — two columns, title uses HARDCODED color (demonstrates the limit of theme swap)
    s = p.slides.add_slide(blank)
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    _set_run(t.text_frame, "What Went Well", size=32, bold=True,
             color=RGBColor(0x1F, 0x4E, 0x79))
    cols_data = [
        ("Customer Wins", "Landed three Fortune 500 accounts; expansions up 22% YoY."),
        ("Product Velocity", "Shipped 14 features this quarter; NPS climbed from 41 to 52."),
    ]
    for ci, (heading, body_txt) in enumerate(cols_data):
        left = Inches(0.75 + ci * 6.3)
        h = s.shapes.add_textbox(left, Inches(1.6), Inches(5.8), Inches(0.6))
        _set_run(h.text_frame, heading, size=22, bold=True, color=NAVY)
        b = s.shapes.add_textbox(left, Inches(2.3), Inches(5.8), Inches(4))
        _set_run(b.text_frame, body_txt, size=18, color=DARK)

    # Slide 5 — closing
    s = p.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, p.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    t = s.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(12.3), Inches(1.5))
    _set_run(t.text_frame, "Thank you", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    INPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    p.save(INPUT_PATH)
    print(f"wrote {INPUT_PATH.relative_to(ROOT)}")


# --- sample_template: "Acme Corporation" branded template -------------------

ACME_ORANGE = RGBColor(0xE8, 0x77, 0x22)
ACME_TEAL = RGBColor(0x00, 0xA9, 0xA5)


def _rewrite_theme(theme_part) -> None:
    """Overwrite theme color scheme and major/minor fonts in-place."""
    root = etree.fromstring(theme_part.blob)
    nsmap = {"a": A_NS}

    color_map = {
        "accent1": "E87722",
        "accent2": "00A9A5",
        "accent3": "4B3F72",
        "accent4": "F2A900",
        "accent5": "6B8E23",
        "accent6": "C8102E",
        "dk2":     "1A1A1A",
        "lt2":     "F7F2E8",
    }
    for name, hex_ in color_map.items():
        for el in root.iter(f"{{{A_NS}}}{name}"):
            for child in list(el):
                el.remove(child)
            etree.SubElement(el, f"{{{A_NS}}}srgbClr", val=hex_)

    for tag, typeface in [("majorFont", "Georgia"), ("minorFont", "Verdana")]:
        for el in root.iter(f"{{{A_NS}}}{tag}"):
            for latin in el.iter(f"{{{A_NS}}}latin"):
                latin.set("typeface", typeface)

    # python-pptx has no public setter for part bytes; _blob is the standard workaround.
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def build_template() -> None:
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)

    master = p.slide_masters[0]
    blank_layout = p.slide_layouts[6]

    # Rewrite theme colors + fonts first.
    theme_part = None
    for rel in master.part.rels.values():
        if rel.reltype.endswith("/theme"):
            theme_part = rel.target_part
            break
    if theme_part is not None:
        _rewrite_theme(theme_part)

    # Branding band + footer need to live on the Blank layout so every slide
    # using Blank inherits them. python-pptx only exposes add_shape on slides,
    # so build the shapes on a scratch slide, then transplant the XML.
    scratch = p.slides.add_slide(blank_layout)

    band = scratch.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, Inches(0.35))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = ACME_ORANGE

    underline = scratch.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.35), p.slide_width, Inches(0.06))
    underline.line.fill.background()
    underline.fill.solid(); underline.fill.fore_color.rgb = ACME_TEAL

    footer = scratch.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(8), Inches(0.35))
    _set_run(footer.text_frame, "ACME CORPORATION  |  CONFIDENTIAL",
             size=10, bold=True, color=ACME_ORANGE)

    layout_sp_tree = blank_layout.shapes._spTree
    for shape in list(scratch.shapes):
        layout_sp_tree.append(deepcopy(shape._element))

    # Properly remove the scratch slide: drop its sldId AND its relationship, so
    # python-pptx doesn't serialize the orphan part (which otherwise produces a
    # duplicate slide1.xml entry in the output zip).
    sld_id_lst = p.slides._sldIdLst
    scratch_sld_id_el = sld_id_lst[-1]
    r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    scratch_rid = scratch_sld_id_el.get(f"{r_ns}id")
    sld_id_lst.remove(scratch_sld_id_el)
    p.part.drop_rel(scratch_rid)

    # Cover slide so the template file is a valid deck and visibly shows the branding.
    cover = p.slides.add_slide(blank_layout)
    t = cover.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(1.2))
    _set_run(t.text_frame, "Acme Corporation Template", size=40, bold=True, color=ACME_TEAL)
    sub = cover.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.6))
    _set_run(sub.text_frame,
             "Replace this deck with a real corporate template, or keep it as a fixture for tests.",
             size=18, color=DARK)

    TEMPLATE_PATH.parent.mkdir(parents=True, exist_ok=True)
    p.save(TEMPLATE_PATH)
    print(f"wrote {TEMPLATE_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    build_input()
    build_template()
